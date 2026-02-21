from pptx import Presentation
from pptx.shapes.placeholder import SlidePlaceholder
from typing import cast
import os


def export_ppt(slides, filename):
    presentation = Presentation()
    slide_layout = presentation.slide_layouts[1]

    for slide_data in slides:
        slide = presentation.slides.add_slide(slide_layout)

        # --- TITLE ---
        title_shape = slide.shapes.title
        if title_shape is not None:
            title = cast(SlidePlaceholder, title_shape)
            title.text = slide_data.title

        # --- CONTENT PLACEHOLDER ---
        content_shape = cast(SlidePlaceholder, slide.placeholders[1])
        text_frame = content_shape.text_frame

        text_frame.clear()
        text_frame.text = slide_data.bullets[0]

        for bullet in slide_data.bullets[1:]:
            p = text_frame.add_paragraph()
            p.text = bullet
            p.level = 1

    os.makedirs("outputs", exist_ok=True)
    file_path = os.path.join("outputs", filename)
    presentation.save(file_path)

    return file_path
